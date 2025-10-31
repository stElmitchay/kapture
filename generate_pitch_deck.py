"""
Generate Kapture pitch deck for hackathon presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    """Title slide with main hook."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color - dark blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # Slate dark

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "KAPTURE"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Get Paid Upfront. Unlock Daily by Working."
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.alignment = PP_ALIGN.CENTER
    tagline_para.font.size = Pt(28)
    tagline_para.font.color.rgb = RGBColor(147, 197, 253)  # Light blue

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Blockchain-Powered Work Verification & Automated Payments"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = RGBColor(203, 213, 225)

def create_problem_slide(prs):
    """The problem slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "The Problem"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 23, 42)

    # Content
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    problems = [
        "Remote workers lack accountability without micromanagement",
        "Employers wait weeks to discover unproductive hires",
        "Freelancers struggle with trust and delayed payments",
        "Traditional time tracking is manual and easily faked",
        "Payment systems require intermediaries and trust"
    ]

    for problem in problems:
        p = tf.add_paragraph()
        p.text = problem
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(20)
        p.level = 0

def create_solution_slide(prs):
    """The solution slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "The Kapture Solution"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Main concept
    concept_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.2))
    concept_frame = concept_box.text_frame
    concept_frame.text = "Flip the payment model: Pay upfront, unlock daily by proving work on the blockchain"
    concept_para = concept_frame.paragraphs[0]
    concept_para.alignment = PP_ALIGN.CENTER
    concept_para.font.size = Pt(26)
    concept_para.font.color.rgb = RGBColor(147, 197, 253)
    concept_para.font.italic = True

    # Features
    features_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(3))
    tf = features_box.text_frame

    features = [
        "Automated work tracking with screenshots, OCR, and activity monitoring",
        "Biometric liveness detection prevents fraud",
        "Blockchain verification for trustless accountability",
        "Instant daily unlocks based on provable work hours",
        "Zero-trust system - no intermediaries needed"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_before = Pt(15)

def create_how_it_works_slide(prs):
    """How it works flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "How It Works"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 23, 42)

    # Steps
    steps_data = [
        ("1", "Employer locks salary in blockchain vault (e.g., $3,000 USDC)"),
        ("2", "Employee runs Kapture - automatic activity tracking begins"),
        ("3", "Work proof generated: screenshots + OCR + biometric verification"),
        ("4", "Daily proof submitted to blockchain oracle"),
        ("5", "Smart contract verifies and unlocks daily payment (e.g., $150)")
    ]

    y_position = 2
    for num, text in steps_data:
        # Number circle
        circle = slide.shapes.add_shape(
            1,  # Circle
            Inches(1.2), Inches(y_position), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(59, 130, 246)  # Blue
        circle.line.color.rgb = RGBColor(59, 130, 246)

        num_text = circle.text_frame
        num_text.text = num
        num_para = num_text.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(20)
        num_para.font.bold = True
        num_para.font.color.rgb = RGBColor(255, 255, 255)

        # Step text
        text_box = slide.shapes.add_textbox(Inches(2), Inches(y_position), Inches(6.5), Inches(0.6))
        text_frame = text_box.text_frame
        text_frame.text = text
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(18)
        text_para.font.color.rgb = RGBColor(51, 65, 85)

        y_position += 0.8

def create_tech_architecture_slide(prs):
    """Technical architecture slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Technical Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Architecture components
    components = [
        ("Desktop Agent (Python)", "Screenshots, OCR (Tesseract), window tracking, biometric detection"),
        ("Oracle Service (Flask API)", "Rate limiting, proof validation, blockchain submission"),
        ("Smart Contract (Solana/Rust)", "Payment vaults, automated unlocks, work verification"),
        ("Storage", "SQLite for local logs, blockchain for immutable records"),
        ("Security", "Face recognition, liveness detection, cryptographic proofs")
    ]

    y_pos = 2
    for title, description in components:
        # Component title
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(7), Inches(0.4))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(147, 197, 253)

        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.35), Inches(7), Inches(0.4))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = RGBColor(203, 213, 225)

        y_pos += 0.9

def create_innovations_slide(prs):
    """Key innovations slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Key Innovations"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 23, 42)

    # Innovations
    innovations_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = innovations_box.text_frame

    innovations = [
        "Reverse Payment Model - Upfront salary locked, unlocked by proof of work",
        "Zero-Trust Verification - Blockchain eliminates need for trust between parties",
        "Anti-Fraud Biometrics - Liveness detection ensures real person working",
        "Automated Accountability - No manual timesheets or tracking needed",
        "Programmable Employment - Smart contracts enforce work terms automatically",
        "Privacy-Preserved Proof - OCR analyzes work without exposing sensitive data"
    ]

    for innovation in innovations:
        p = tf.add_paragraph()
        p.text = innovation
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(18)
        p.font.bold = True

        # Extract description after dash
        if " - " in innovation:
            parts = innovation.split(" - ", 1)
            p.text = parts[0]
            p.font.color.rgb = RGBColor(59, 130, 246)

            # Add description as sub-paragraph
            p2 = tf.add_paragraph()
            p2.text = parts[1]
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(71, 85, 105)
            p2.space_before = Pt(3)
            p2.level = 1

def create_use_cases_slide(prs):
    """Use cases slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Use Cases"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Use cases in grid
    use_cases = [
        ("Remote Teams", "Accountability without micromanagement"),
        ("DAO Contributors", "Trustless payment for decentralized work"),
        ("Freelancers", "Build trust with new clients instantly"),
        ("Bootcamps", "Students unlock refunds by completing work"),
        ("Self-Discipline", "Lock your own money to force productivity"),
        ("Global Workforce", "Borderless payments with verification")
    ]

    # Grid layout: 2 columns, 3 rows
    x_positions = [1.2, 5.2]
    y_positions = [2, 3.4, 4.8]

    idx = 0
    for y in y_positions:
        for x in x_positions:
            if idx >= len(use_cases):
                break
            title, desc = use_cases[idx]

            # Box
            box = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(x), Inches(y), Inches(3.5), Inches(1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(30, 41, 59)
            box.line.color.rgb = RGBColor(59, 130, 246)
            box.line.width = Pt(2)

            # Title
            text_frame = box.text_frame
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)

            p1 = text_frame.paragraphs[0]
            p1.text = title
            p1.font.size = Pt(18)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(147, 197, 253)

            # Description
            p2 = text_frame.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(13)
            p2.font.color.rgb = RGBColor(203, 213, 225)
            p2.space_before = Pt(5)

            idx += 1

def create_demo_slide(prs):
    """Demo/current state slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Live & Functional"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 23, 42)

    # Current status
    status_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = status_box.text_frame

    # Main status
    p1 = tf.paragraphs[0]
    p1.text = "✓ Published on PyPI - pip3 install loggerheads"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(22, 163, 74)  # Green

    # Features completed
    features = [
        "Desktop tracking agent with OCR and biometrics",
        "Oracle API with rate limiting and validation",
        "Solana smart contracts deployed",
        "CLI interface with employer/employee workflows",
        "Automated daily submissions and unlocks",
        "SQLite logging and blockchain persistence"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = f"✓ {feature}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(12)

    # Call to action
    cta_para = tf.add_paragraph()
    cta_para.text = "Try it now: pip3 install loggerheads"
    cta_para.font.size = Pt(20)
    cta_para.font.bold = True
    cta_para.font.color.rgb = RGBColor(59, 130, 246)
    cta_para.space_before = Pt(25)
    cta_para.alignment = PP_ALIGN.CENTER

def create_market_slide(prs):
    """Market opportunity slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Market Opportunity"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Market data
    markets_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = markets_box.text_frame

    markets = [
        ("1.57B", "Remote workers worldwide by 2025"),
        ("$9.2T", "Global freelance market size"),
        ("20,000+", "Active DAOs managing contributor payments"),
        ("$160B", "Time tracking & workforce management software market"),
        ("Growing", "Blockchain adoption in HR & payroll systems")
    ]

    for stat, description in markets:
        # Stat
        p1 = tf.add_paragraph()
        p1.text = stat
        p1.font.size = Pt(32)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(147, 197, 253)
        p1.space_before = Pt(15)

        # Description
        p2 = tf.add_paragraph()
        p2.text = description
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(203, 213, 225)
        p2.space_before = Pt(3)
        p2.space_after = Pt(10)

def create_vision_slide(prs):
    """Vision and next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Vision & Roadmap"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 23, 42)

    # Vision statement
    vision_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(1))
    vision_frame = vision_box.text_frame
    vision_frame.text = "Make work verification trustless, automatic, and borderless"
    vision_para = vision_frame.paragraphs[0]
    vision_para.alignment = PP_ALIGN.CENTER
    vision_para.font.size = Pt(24)
    vision_para.font.color.rgb = RGBColor(59, 130, 246)
    vision_para.font.italic = True

    # Next steps
    steps_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(3))
    tf = steps_box.text_frame

    p_title = tf.paragraphs[0]
    p_title.text = "Next Steps:"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(15, 23, 42)

    next_steps = [
        "Cross-platform support (Windows, Linux desktop apps)",
        "Mobile companion app for on-the-go workers",
        "Multi-chain support (Ethereum, Polygon, Base)",
        "Integration with popular HR platforms (BambooHR, Workday)",
        "AI-powered work quality analysis (not just hours)",
        "Marketplace for verified workers and employers"
    ]

    for step in next_steps:
        p = tf.add_paragraph()
        p.text = f"→ {step}"
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(10)

def create_thank_you_slide(prs):
    """Final thank you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)

    # Main message
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "KAPTURE"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Work Verified. Payments Unlocked. Trust Eliminated."
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.alignment = PP_ALIGN.CENTER
    tagline_para.font.size = Pt(26)
    tagline_para.font.color.rgb = RGBColor(147, 197, 253)

    # Links
    links_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    links_frame = links_box.text_frame

    p1 = links_frame.paragraphs[0]
    p1.text = "GitHub: github.com/stElmitchay/loggerheads"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(18)
    p1.font.color.rgb = RGBColor(203, 213, 225)

    p2 = links_frame.add_paragraph()
    p2.text = "PyPI: pip3 install loggerheads"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_before = Pt(10)

def main():
    """Generate the complete pitch deck."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Generating pitch deck slides...")

    create_title_slide(prs)
    print("✓ Title slide")

    create_problem_slide(prs)
    print("✓ Problem slide")

    create_solution_slide(prs)
    print("✓ Solution slide")

    create_how_it_works_slide(prs)
    print("✓ How It Works slide")

    create_tech_architecture_slide(prs)
    print("✓ Technical Architecture slide")

    create_innovations_slide(prs)
    print("✓ Key Innovations slide")

    create_use_cases_slide(prs)
    print("✓ Use Cases slide")

    create_demo_slide(prs)
    print("✓ Demo/Live Product slide")

    create_market_slide(prs)
    print("✓ Market Opportunity slide")

    create_vision_slide(prs)
    print("✓ Vision & Roadmap slide")

    create_thank_you_slide(prs)
    print("✓ Thank You slide")

    # Save presentation
    filename = "Kapture_Pitch_Deck.pptx"
    prs.save(filename)
    print(f"\n✅ Pitch deck generated: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
